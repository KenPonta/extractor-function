from io import BytesIO
from types import SimpleNamespace
from unittest import mock

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

import placeholder_extractor as px


# --- fixtures --------------------------------------------------------------- #
def _png_bytes(color="red"):
    buf = BytesIO()
    Image.new("RGB", (12, 12), color).save(buf, format="PNG")
    return buf.getvalue()


def _make_deck(path, images=0):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Hello world"
    for _ in range(images):
        img = path.parent / "pic.png"
        img.write_bytes(_png_bytes())
        slide.shapes.add_picture(str(img), Inches(1), Inches(2), Inches(2), Inches(2))
    prs.save(str(path))
    return path


@pytest.fixture
def deck(tmp_path):
    return _make_deck(tmp_path / "deck.pptx", images=1)


@pytest.fixture
def fake_openai():
    # Patch openai.OpenAI so no real request is made; responses.create returns a canned reply.
    client = mock.MagicMock()
    client.responses.create.return_value = SimpleNamespace(output_text="A red square.")
    with mock.patch("openai.OpenAI", return_value=client):
        yield client


# --- pure helpers ----------------------------------------------------------- #
def test_xml_attr_escapes():
    assert px._xml_attr('a & b < c > "d"') == "a &amp; b &lt; c &gt; &quot;d&quot;"


def test_data_url_png():
    assert px._data_url(b"\x89PNG", "png").startswith("data:image/png;base64,")


def test_data_url_jpg_maps_to_jpeg():
    assert px._data_url(b"\xff\xd8", "jpg").startswith("data:image/jpeg;base64,")


def test_data_url_unknown_ext_falls_back_to_png():
    assert px._data_url(b"x", "bmp").startswith("data:image/png;base64,")


def test_sort_key_handles_none():
    assert px._sort_key(SimpleNamespace(top=None, left=5)) == (0, 5)


# --- validation ------------------------------------------------------------- #
def test_validate_missing(tmp_path):
    with pytest.raises(FileNotFoundError):
        px._validate_pptx(tmp_path / "nope.pptx")


def test_validate_wrong_extension(tmp_path):
    bad = tmp_path / "a.txt"
    bad.write_text("x")
    with pytest.raises(ValueError, match="file type mismatch"):
        px._validate_pptx(bad)


def test_validate_ok(deck):
    assert px._validate_pptx(deck) == deck


# --- data model ------------------------------------------------------------- #
def test_config_defaults_and_override():
    assert px.ExtractorConfig().model == "gpt-4.1"
    assert px.ExtractorConfig(model="gpt-4o").model == "gpt-4o"


def test_slides_do_not_share_block_list():
    a, b = px.Slide(1), px.Slide(2)
    a.blocks.append(("text", "x"))
    assert b.blocks == []


# --- rendering (pure, no network) ------------------------------------------- #
def test_render_xml_structure():
    slides = [px.Slide(1, [("text", "Title"), ("image", 1)])]
    by_id = {1: px.ImageRef(1, b"", "png", "A chart.")}
    xml = px.PlaceholderExtractor()._render_xml(slides, by_id, "deck.pptx")
    assert "<documents>" in xml
    assert '<document index="1" filename="deck.pptx" data-type="PPTX">' in xml
    assert "Slide 1" in xml
    assert "[Image 1]\nA chart." in xml
    assert xml.strip().endswith("</documents>")


def test_render_xml_escapes_filename():
    xml = px.PlaceholderExtractor()._render_xml([], {}, "a & b.pptx")
    assert 'filename="a &amp; b.pptx"' in xml


# --- end to end with mocked OpenAI ------------------------------------------ #
def test_extract_xml_end_to_end(deck, fake_openai):
    xml = px.PlaceholderExtractor().extract_xml(deck)
    assert "Hello world" in xml                              # native text, no model
    assert "A red square." in xml                            # mocked image description
    assert fake_openai.responses.create.call_count == 1      # one image -> one call


def test_image_is_deduplicated(tmp_path, fake_openai):
    deck = _make_deck(tmp_path / "dup.pptx", images=2)       # same image twice
    px.PlaceholderExtractor().extract_xml(deck)
    assert fake_openai.responses.create.call_count == 1      # described once
