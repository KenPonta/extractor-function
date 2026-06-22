from pptx import Presentation

prs = Presentation("/Users/pontakornkangvanwanich/extaction-function/Panel TFP Project Slide.pptx")

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(shape.text_frame.text)