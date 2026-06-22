#!/usr/bin/env python3
"""
pptx_extract.py
===============
Extract a PowerPoint deck slide-by-slide.

Each slide is routed independently by a cheap, render-free heuristic:
  * text-heavy slides  -> pulled out directly with python-pptx (free, fast)
  * visual slides       -> rendered to an image and read by an OpenAI GPT vision model

Only the slides that actually need vision are rendered and sent to the model, so
runtime/cost scales with the number of *visual* slides, not the whole deck.

Output is a Markdown document with one `## Slide N (text|vision)` section per slide,
in original order.

Dependencies
------------
    pip install python-pptx pymupdf openai python-dotenv
    # plus LibreOffice on the system (provides `soffice`) for rendering:
    #   Debian/Ubuntu:  sudo apt-get install libreoffice
    #   macOS (brew):   brew install --cask libreoffice

Setup
-----
    # Create a .env file in the same directory with:
    # OPENAI_API_KEY=sk-...
    # (only needed if any slide routes to vision)

Usage
-----
    python pptx_extract.py deck.pptx                 # -> deck.md
    python pptx_extract.py deck.pptx -o out.md
    python pptx_extract.py deck.pptx -m gpt-5.5 --dpi 150
"""

from __future__ import annotations

import argparse
import base64
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from dotenv import load_dotenv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Load environment variables from .env file
load_dotenv()


# --------------------------------------------------------------------------- #
# Tunable configuration
# --------------------------------------------------------------------------- #
MODEL = "gpt-4.1"          # any vision-capable OpenAI model (gpt-5.5, gpt-4.1, gpt-4o, ...)
RENDER_DPI = 150           # rasterization quality for vision slides
IMAGE_DETAIL = "high"      # "low" | "high" | "auto"  (high = better on dense slides, costs more)

# Routing thresholds (see route_slide). These are heuristics — tune to your decks.
WORD_CAP = 60              # per-shape word cap so one wall-of-text box can't dominate
VISION_PIC_FRACTION = 0.45 # slide area covered by pictures above which it leans "visual"
VISION_TEXT_BELOW = 30     # weighted text score below which a picture-heavy slide -> vision
SPARSE_TEXT = 5            # near-empty text score
SPARSE_PIC_FRACTION = 0.10 # any non-trivial picture presence

VISION_PROMPT = (
    "You are reading a single slide from a presentation. Transcribe everything on it "
    "as clean Markdown: every piece of text (title, bullets, labels, captions, footnotes), "
    "and for any chart, diagram, table, screenshot, or image, describe what it shows and "
    "report the data or relationships it conveys (axis values, trends, comparisons, flow). "
    "Preserve the reading order. Do not add commentary and do not invent anything that is "
    "not visible on the slide."
)


# --------------------------------------------------------------------------- #
# Text weighting (used by the router)
# --------------------------------------------------------------------------- #
def text_weight_for(shape) -> float:
    """Weight a text container by where it lives. Titles/body count most;
    footers/page-numbers/dates are boilerplate and count nothing."""
    if not getattr(shape, "is_placeholder", False):
        return 0.8                                  # loose text box / callout
    t = str(getattr(shape.placeholder_format, "type", "") or "")
    if "SUBTITLE" in t:                             # check before TITLE (substring overlap)
        return 1.2
    if "TITLE" in t:                                # TITLE, CENTER_TITLE
        return 1.5
    if "BODY" in t or "OBJECT" in t:
        return 1.0
    if any(k in t for k in ("FOOTER", "SLIDE_NUMBER", "DATE")):
        return 0.0                                  # boilerplate -> ignore
    return 0.8


def text_score(shape) -> float:
    """Weighted, capped word count for one text-bearing shape."""
    words = len(shape.text_frame.text.split())
    return text_weight_for(shape) * min(words, WORD_CAP)


# --------------------------------------------------------------------------- #
# Per-slide routing  (render-free, microseconds per slide)
# --------------------------------------------------------------------------- #
def route_slide(slide, slide_area: int) -> str:
    """Return "text" or "vision" for a single slide.

    Native PowerPoint charts and tables are readable directly, so they do NOT
    push a slide toward vision — only pasted-image content (pictures) does.
    """
    score = 0.0
    pic_area = 0

    def walk(shapes):
        nonlocal score, pic_area
        for s in shapes:
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(s.shapes)
            elif s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pic_area += (s.width or 0) * (s.height or 0)
            elif getattr(s, "has_text_frame", False):
                score += text_score(s)
            # native chart/table -> handled by extract_text, no vision needed

    walk(slide.shapes)
    pic_fraction = pic_area / slide_area if slide_area else 0.0

    if pic_fraction > VISION_PIC_FRACTION and score < VISION_TEXT_BELOW:
        return "vision"                 # dominated by imagery, little text
    if score < SPARSE_TEXT and pic_fraction > SPARSE_PIC_FRACTION:
        return "vision"                 # near-empty text; content lives in the picture
    return "text"                       # title+bullets, tables, native charts extract fine


# --------------------------------------------------------------------------- #
# Cheap path: direct text extraction with python-pptx
# --------------------------------------------------------------------------- #
def extract_text(slide) -> str:
    """Pull text frames, table cells, and speaker notes from a slide."""
    chunks: list[str] = []

    def walk(shapes):
        for s in shapes:
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(s.shapes)
                continue
            if getattr(s, "has_text_frame", False):
                txt = s.text_frame.text.strip()
                if txt:
                    chunks.append(txt)
            if getattr(s, "has_table", False):
                for row in s.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        chunks.append(" | ".join(cells))

    walk(slide.shapes)

    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            chunks.append("**Speaker notes:** " + notes)

    return "\n\n".join(chunks)


# --------------------------------------------------------------------------- #
# Vision path: render flagged slides, then read them with GPT
# --------------------------------------------------------------------------- #
def render_pptx_to_pdf(pptx_path: Path, workdir: Path) -> Path:
    """Convert the whole deck to a PDF once, using headless LibreOffice."""
    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not soffice:
        raise RuntimeError(
            "LibreOffice ('soffice') not found. Install it to render visual slides "
            "(e.g. `sudo apt-get install libreoffice`)."
        )
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(workdir), str(pptx_path)],
        check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
    )
    pdf_path = workdir / (pptx_path.stem + ".pdf")
    if not pdf_path.exists():
        raise RuntimeError("LibreOffice did not produce a PDF.")
    return pdf_path


def render_pages_to_pngs(pdf_path: Path, slide_numbers: list[int],
                         workdir: Path, dpi: int = RENDER_DPI) -> dict[int, Path]:
    """Rasterize ONLY the requested pages (1-based slide numbers) to PNG."""
    import fitz  # PyMuPDF
    fitz.TOOLS.mupdf_display_errors(False)  # benign LibreOffice-PDF warnings -> quiet

    out: dict[int, Path] = {}
    if not slide_numbers:
        return out
    doc = fitz.open(str(pdf_path))
    try:
        for n in slide_numbers:
            idx = n - 1                              # slide N -> PDF page N-1 (1 page/slide)
            if idx < 0 or idx >= doc.page_count:
                print(f"  ! slide {n}: no matching PDF page, skipping", file=sys.stderr)
                continue
            pix = doc[idx].get_pixmap(dpi=dpi)
            png = workdir / f"slide_{n}.png"
            pix.save(str(png))
            out[n] = png
    finally:
        doc.close()
    return out


def describe_slide_with_gpt(client, png_path: Path,
                            model: str = MODEL, detail: str = IMAGE_DETAIL) -> str:
    """Send one rendered slide image to an OpenAI GPT vision model (Responses API)."""
    b64 = base64.b64encode(png_path.read_bytes()).decode("utf-8")
    resp = client.responses.create(
        model=model,
        input=[{
            "role": "user",
            "content": [
                {"type": "input_text", "text": VISION_PROMPT},
                {"type": "input_image",
                 "image_url": f"data:image/png;base64,{b64}",
                 "detail": detail},
            ],
        }],
    )
    return (resp.output_text or "").strip()


# --------------------------------------------------------------------------- #
# Orchestration
# --------------------------------------------------------------------------- #
def process_pptx(pptx_path: Path, out_path: Path | None = None,
                 model: str = MODEL, dpi: int = RENDER_DPI) -> str:
    prs = Presentation(str(pptx_path))
    area = (prs.slide_width or 1) * (prs.slide_height or 1)

    # 1) Route every slide independently; grab text now for the cheap path.
    routed: list[tuple[int, str, str | None]] = []
    for i, slide in enumerate(prs.slides, start=1):
        kind = route_slide(slide, area)
        routed.append((i, kind, extract_text(slide) if kind == "text" else None))

    vision_pages = [i for i, kind, _ in routed if kind == "vision"]
    print(f"{len(routed)} slides: {len(routed) - len(vision_pages)} text, "
          f"{len(vision_pages)} vision {vision_pages or ''}", file=sys.stderr)

    # 2) Render + GPT only the flagged slides.
    vision_text: dict[int, str] = {}
    if vision_pages:
        if not os.environ.get("OPENAI_API_KEY"):
            sys.exit("OPENAI_API_KEY is not set, but some slides need the vision model.")
        from openai import OpenAI
        client = OpenAI()
        with tempfile.TemporaryDirectory() as tmp:
            workdir = Path(tmp)
            pdf_path = render_pptx_to_pdf(pptx_path, workdir)
            pngs = render_pages_to_pngs(pdf_path, vision_pages, workdir, dpi=dpi)
            for n in vision_pages:
                if n not in pngs:
                    vision_text[n] = "_[slide could not be rendered]_"
                    continue
                print(f"  vision: slide {n} -> {model}", file=sys.stderr)
                try:
                    vision_text[n] = describe_slide_with_gpt(client, pngs[n], model=model)
                except Exception as exc:                       # keep going on per-slide failures
                    vision_text[n] = f"_[vision call failed: {exc}]_"

    # 3) Stitch back together in slide order.
    blocks: list[str] = []
    for i, kind, text in routed:
        body = text if kind == "text" else vision_text.get(i, "")
        body = body.strip() if body else "_[no content]_"
        blocks.append(f"## Slide {i}  ({kind})\n\n{body}")
    result = "\n\n---\n\n".join(blocks) + "\n"

    if out_path:
        out_path.write_text(result, encoding="utf-8")
    return result


def main() -> None:
    ap = argparse.ArgumentParser(
        description="Extract a PowerPoint slide-by-slide, routing visual slides to a GPT vision model."
    )
    ap.add_argument("pptx", help="path to the .pptx file")
    ap.add_argument("-o", "--out", help="output .md path (default: <pptx name>.md)")
    ap.add_argument("-m", "--model", default=MODEL, help=f"OpenAI vision model (default: {MODEL})")
    ap.add_argument("--dpi", type=int, default=RENDER_DPI, help=f"render DPI (default: {RENDER_DPI})")
    args = ap.parse_args()

    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        sys.exit(f"File not found: {pptx_path}")
    out_path = Path(args.out) if args.out else pptx_path.with_suffix(".md")

    process_pptx(pptx_path, out_path, model=args.model, dpi=args.dpi)
    print(f"[written to {out_path}]", file=sys.stderr)


if __name__ == "__main__":
    main()